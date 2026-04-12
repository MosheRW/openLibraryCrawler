"""
One-shot script: generates presentation.pptx from the slide content.
Run: python generate_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x1E, 0x1E, 0x2E)   # near-black navy
ACCENT       = RGBColor(0x74, 0xC7, 0xEC)   # sky blue
ACCENT2      = RGBColor(0xA6, 0xE3, 0xA1)   # mint green
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xDD)
CODE_BG      = RGBColor(0x18, 0x18, 0x28)
CODE_FG      = RGBColor(0xA6, 0xE3, 0xA1)
TABLE_HEADER = RGBColor(0x31, 0x32, 0x44)
TABLE_ROW1   = RGBColor(0x28, 0x29, 0x3A)
TABLE_ROW2   = RGBColor(0x22, 0x23, 0x32)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


def add_slide():
    return prs.slides.add_slide(blank_layout)


def set_bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_title_slide():
    slide = add_slide()
    set_bg(slide)

    # accent bar top
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.07))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # main title
    add_textbox(slide,
                "OpenLibrary Crawler",
                Inches(1), Inches(1.8), Inches(11.3), Inches(1.4),
                font_size=48, bold=True, color=ACCENT,
                align=PP_ALIGN.CENTER)

    # subtitle
    add_textbox(slide,
                "Web Automation & Performance Testing\nwith Python + Playwright",
                Inches(1), Inches(3.3), Inches(11.3), Inches(1.2),
                font_size=24, color=LIGHT_GRAY,
                align=PP_ALIGN.CENTER)

    # bottom accent bar
    bar2 = slide.shapes.add_shape(1, 0, SLIDE_H - Inches(0.07), SLIDE_W, Inches(0.07))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = ACCENT
    bar2.line.fill.background()


def add_section_slide(title, bullets, bullet_color=LIGHT_GRAY,
                      highlight_first=False):
    """Generic title + bullet list slide."""
    slide = add_slide()
    set_bg(slide)

    # title bar
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = TABLE_HEADER
    bar.line.fill.background()

    add_textbox(slide, title,
                Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.85),
                font_size=28, bold=True, color=ACCENT)

    # bullets
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.25),
                                     Inches(12.3), Inches(6.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        indent = 0
        text = line
        if line.startswith("    "):
            indent = 1
            text = line.lstrip()
        elif line.startswith("  "):
            indent = 1
            text = line.lstrip()

        p.level = indent
        run = p.add_run()
        run.text = text
        run.font.size = Pt(16 if indent == 0 else 14)
        run.font.color.rgb = (ACCENT if (highlight_first and i == 0)
                               else bullet_color)

    return slide


def add_table_slide(title, headers, rows, col_widths=None):
    slide = add_slide()
    set_bg(slide)

    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = TABLE_HEADER
    bar.line.fill.background()

    add_textbox(slide, title,
                Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.85),
                font_size=28, bold=True, color=ACCENT)

    ncols = len(headers)
    nrows = len(rows) + 1
    if col_widths is None:
        col_widths = [SLIDE_W / ncols] * ncols

    table_w = sum(col_widths)
    left = (SLIDE_W - table_w) / 2
    top  = Inches(1.25)
    row_h = min(Inches(0.55), (SLIDE_H - top - Inches(0.2)) / nrows)

    tbl = slide.shapes.add_table(
        nrows, ncols, left, top,
        int(table_w), int(row_h * nrows)
    ).table

    for c, w in enumerate(col_widths):
        tbl.columns[c].width = int(w)

    # header row
    for c, hdr in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = hdr
        run.font.bold = True
        run.font.size = Pt(13)
        run.font.color.rgb = BG_DARK
        p.alignment = PP_ALIGN.CENTER

    # data rows
    for r, row in enumerate(rows):
        bg = TABLE_ROW1 if r % 2 == 0 else TABLE_ROW2
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(12)
            run.font.color.rgb = WHITE

    return slide


def add_code_slide(title, code_text, bullets_below=None):
    slide = add_slide()
    set_bg(slide)

    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = TABLE_HEADER
    bar.line.fill.background()

    add_textbox(slide, title,
                Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.85),
                font_size=28, bold=True, color=ACCENT)

    code_h = Inches(4.5) if bullets_below else Inches(5.8)
    code_box = slide.shapes.add_shape(1,
                                      Inches(0.4), Inches(1.2),
                                      Inches(12.5), code_h)
    code_box.fill.solid(); code_box.fill.fore_color.rgb = CODE_BG
    code_box.line.color.rgb = ACCENT

    add_textbox(slide, code_text,
                Inches(0.55), Inches(1.3),
                Inches(12.2), code_h - Inches(0.2),
                font_size=12, color=CODE_FG)

    if bullets_below:
        txBox = slide.shapes.add_textbox(Inches(0.5),
                                         Inches(1.25) + code_h,
                                         Inches(12.3),
                                         SLIDE_H - Inches(1.25) - code_h - Inches(0.1))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(bullets_below):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(14)
            run.font.color.rgb = LIGHT_GRAY


# ── SLIDE 1: Title ───────────────────────────────────────────────────────────
add_title_slide()

# ── SLIDE 2: What Is OpenLibrary? ────────────────────────────────────────────
add_section_slide(
    "What Is OpenLibrary?",
    [
        "Free, open-source digital library run by the Internet Archive",
        "25 million+ book records — community-edited, wiki-style",
        "Unique among book sites:",
        "  • Free digital lending (1-hour loans)",
        "  • Public API — no key required for reads",
        "  • Fully open-source codebase on GitHub",
        "  • Open data dumps (CC0 license)",
        "",
        "Reading-list UI (Want to Read / Currently Reading / Already Read)",
        "is not fully covered by the API alone — an interesting automation target",
    ]
)

# ── SLIDE 3: Project Goal ────────────────────────────────────────────────────
add_section_slide(
    "Project Goal",
    [
        "Automate the full book-management lifecycle and measure real site performance:",
        "",
        "  1.  Authenticate as a real user",
        "  2.  Search books by title + publication year",
        "  3.  Add books to reading shelves",
        "  4.  Assert shelf counts match expectations",
        "  5.  Measure page load performance (First Paint, DCL, Load Time)",
        "  6.  Generate an interactive HTML report",
        "",
        "Tech stack:  Python 3.10+  ·  Playwright (async)  ·  asyncio  ·  PyYAML  ·  Chart.js",
    ]
)

# ── SLIDE 4: Architecture ────────────────────────────────────────────────────
add_code_slide(
    "Architecture — Page Object Model",
    """\
main.py
  └── orchestrator()
        ├── ProfilePage       ← shelf management & assertions
        ├── HomePage          ← search input
        ├── SearchResultsPage ← result parsing + pagination
        └── BookPage          ← "Add to shelf" actions

All pages inherit BasePage → login via AuthPage""",
    bullets_below=[
        "Pattern: Page Object Model — one class per page, selectors isolated inside it",
        "Pattern: Singleton — Browser, Logger, Config, Results (one instance per run)",
        "Pattern: Factory Method — PageClass.create() handles async initialisation",
    ]
)

# ── SLIDE 5: Core Workflow ───────────────────────────────────────────────────
add_code_slide(
    "Core Workflow",
    """\
options.yaml
  ↓
Config (singleton) — credentials + queries + settings
  ↓
ProfilePage.create() → login → navigate to profile
  ↓  (optional) clear shelves
  ↓
For each Query:
    search_books_by_title_under_year()  → list of URLs
    add_books_to_reading_list(urls)     → (want_count, read_count)
    assert_reading_list_count(expected) → pass / AssertionError
  ↓
logger.save_logs()  →  performance_report.json
generate_report()   →  report.html""",
    bullets_below=[
        "Shelf counts accumulate across queries — each assertion checks total shelf state, not just the current query's books",
    ]
)

# ── SLIDE 6: Performance Monitoring ─────────────────────────────────────────
add_table_slide(
    "Performance Monitoring",
    ["Metric", "Source", "Threshold"],
    [
        ["first_paint_ms",        "performance.getEntriesByType('paint')",                      "—"],
        ["dom_content_loaded_ms", "PerformanceNavigationTiming.domContentLoadedEventEnd",        "—"],
        ["load_time_ms",          "PerformanceNavigationTiming.loadEventEnd",                    "configurable"],
        ["is_within_threshold",   "load_time_ms ≤ threshold",                                   "BookPage 2500 ms"],
        ["",                      "",                                                             "ProfilePage 2000 ms"],
        ["",                      "",                                                             "SearchResults 5000 ms"],
    ],
    col_widths=[Inches(2.8), Inches(6.5), Inches(3.5)]
)

# ── SLIDE 7: Key Design Decisions ───────────────────────────────────────────
add_table_slide(
    "Key Design Decisions",
    ["Decision", "Rationale"],
    [
        ["Playwright over Selenium",       "Native async, built-in auto-wait, faster, modern API"],
        ["Singleton Browser",              "One tab reused — no context-switching overhead"],
        ["create() factory on each page",  "async __init__ is invalid in Python; factory encapsulates setup"],
        ["Book with 6 export formats",     "JSON / XML / CSV / YAML / HTML / Markdown for downstream use"],
        ["YAML config (not CLI args)",     "Nested queries + credentials are cleaner than flags"],
        ["ReadingStatus Enum",             "Replaces fragile bool | str return type with explicit states"],
    ],
    col_widths=[Inches(4.0), Inches(8.8)]
)

# ── SLIDE 8: Engineering Challenges ─────────────────────────────────────────
add_section_slide(
    "Interesting Engineering Challenges",
    [
        "1.  Stale UI during bulk shelf clearing",
        "    OpenLibrary's remove button stops responding after ~5 removals.",
        "    Fix: reload every ITERATIONS_BEFORE_RELOAD operations to reset the DOM.",
        "",
        "2.  Shelf button has three states",
        "    Unset / set to this shelf / set to a different shelf.",
        "    Modelled with ReadingStatus Enum  (SUCCESS=1, FAILURE=0, NOT_SUCCESS=-1).",
        "",
        "3.  No reliable post-login DOM signal",
        "    Current: wait_for_timeout(5000).  Correct fix: wait_for_url().",
        "",
        "4.  Cumulative shelf counts across queries",
        "    Counters carry over; each assertion checks total shelf state.",
        "",
        "5.  Duplicate performance log entries",
        "    navigate() and is_book_marked_as() both log the same URL.",
        "    Fix: deduplicate by (url, page) key in report generator.",
    ]
)

# ── SLIDE 9: OpenLibrary vs. Similar Sites ───────────────────────────────────
add_table_slide(
    "OpenLibrary vs. Similar Sites",
    ["Feature", "OpenLibrary", "Goodreads", "LibraryThing", "Google Books"],
    [
        ["Owner",            "Internet Archive",  "Amazon",          "TinyCat",        "Google"],
        ["API",              "Free, no key",      "Deprecated 2020", "Key required",   "Key + quota"],
        ["Open source",      "Yes",               "No",              "Partial",        "No"],
        ["Digital lending",  "Yes",               "No",              "No",             "Preview only"],
        ["Community editing","Yes (wiki)",         "No",              "Yes",            "No"],
        ["Open data",        "Yes (CC0 dump)",     "No",              "No",             "No"],
        ["Automatable?",     "Yes — UI + API",    "Very limited",    "Yes (key)",      "Yes (quota)"],
    ],
    col_widths=[Inches(2.2), Inches(2.5), Inches(2.2), Inches(2.5), Inches(2.5)]
)

# ── SLIDE 10: Future Improvements ────────────────────────────────────────────
add_section_slide(
    "Future Improvements",
    [
        "1.  Replace wait_for_timeout(5000) post-login with wait_for_url()",
        "    Eliminates network-dependent flakiness",
        "",
        "2.  Add shelf field to Query config",
        "    Removes random shelf assignment — makes runs reproducible",
        "",
        "3.  Parallel query execution with asyncio.gather() + Page pool",
        "    Each query worker gets its own page; Logger/Results singletons already safe",
        "",
        "4.  Retry decorator for flaky navigation & book-add operations",
        "",
        "5.  Fix report gallery — glob('*.png') → glob('**/*.png')",
        "    Book-page screenshots in subdirectories currently missing from the report",
        "",
        "6.  Persist results to SQLite — track performance trends across runs",
        "",
        "7.  Docker + CI — run on every commit, detect regressions automatically",
    ]
)

out_path = "presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
